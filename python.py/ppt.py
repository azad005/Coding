from pptx import Presentation

# Create a presentation object
presentation = Presentation()

# Slide 1: Title Slide
slide = presentation.slides.add_slide(presentation.slide_layouts[0])
title = slide.shapes.title
subtitle = slide.placeholders[1]
title.text = "Music Hub"
subtitle.text = "SDLC Implementation with 7 Stages\nPresented by: [Your Name]\nDate: [Insert Date]"

# Slide 2: Introduction
slide = presentation.slides.add_slide(presentation.slide_layouts[1])
title = slide.shapes.title
content = slide.placeholders[1]
title.text = "Introduction"
content.text = (
    "Objective:\n- Create an all-in-one music platform for streaming, managing playlists, and user interaction.\n\n"
    "Purpose:\n- Deliver a user-friendly and engaging music experience.\n\n"
    "Scope:\n- Available on mobile and desktop."
)

# Slide 3: SDLC Stages Overview
slide = presentation.slides.add_slide(presentation.slide_layouts[1])
title = slide.shapes.title
content = slide.placeholders[1]
title.text = "SDLC Stages Overview"
content.text = (
    "Stages:\n1. Planning\n2. Requirement Analysis\n3. System Design\n4. Coding\n"
    "5. Testing\n6. Deployment\n7. Maintenance\n\n"
    "Model Used: [Waterfall/Iterative/Agile]"
)

# Slide 4: Planning
slide = presentation.slides.add_slide(presentation.slide_layouts[1])
title = slide.shapes.title
content = slide.placeholders[1]
title.text = "Planning"
content.text = (
    "Goals:\n- Define project objectives.\n- Set timeline and budget.\n\n"
    "Deliverables:\n- Feasibility report.\n- Project charter."
)

# Slide 5: Requirement Analysis
slide = presentation.slides.add_slide(presentation.slide_layouts[1])
title = slide.shapes.title
content = slide.placeholders[1]
title.text = "Requirement Analysis"
content.text = (
    "Activities:\n- Gather functional and non-functional requirements.\n- Consult stakeholders.\n\n"
    "Outcome:\n- Use case diagrams.\n- Requirement specification document."
)

# Slide 6: System Design
slide = presentation.slides.add_slide(presentation.slide_layouts[1])
title = slide.shapes.title
content = slide.placeholders[1]
title.text = "System Design"
content.text = (
    "Focus:\n- High-level and detailed design.\n- Database schema and user interface design.\n\n"
    "Tools:\n- UML diagrams, flowcharts, ERD."
)

# Slide 7: Coding
slide = presentation.slides.add_slide(presentation.slide_layouts[1])
title = slide.shapes.title
content = slide.placeholders[1]
title.text = "Coding"
content.text = (
    "Activities:\n- Develop front-end (e.g., React/Flutter).\n- Develop back-end (e.g., Node.js, Python).\n"
    "- Integrate database (e.g., MySQL/MongoDB).\n\n"
    "Languages:\n- Mention technologies and frameworks.\n\n"
    "Outcome:\n- Functional modules."
)

# Slide 8: Testing
slide = presentation.slides.add_slide(presentation.slide_layouts[1])
title = slide.shapes.title
content = slide.placeholders[1]
title.text = "Testing"
content.text = (
    "Types of Testing:\n- Unit Testing.\n- Integration Testing.\n- System Testing.\n- User Acceptance Testing (UAT).\n\n"
    "Tools:\n- Selenium, JUnit, etc."
)

# Slide 9: Deployment
slide = presentation.slides.add_slide(presentation.slide_layouts[1])
title = slide.shapes.title
content = slide.placeholders[1]
title.text = "Deployment"
content.text = (
    "Environment:\n- Production server setup.\n- Deployment tools (e.g., Docker, Jenkins).\n\n"
    "Activities:\n- Deploy the application.\n- Provide initial training to users."
)

# Slide 10: Maintenance
slide = presentation.slides.add_slide(presentation.slide_layouts[1])
title = slide.shapes.title
content = slide.placeholders[1]
title.text = "Maintenance"
content.text = (
    "Post-Deployment Activities:\n- Bug fixes and updates.\n- Adding new features based on user feedback.\n\n"
    "Support:\n- Regular monitoring and backups."
)

# Slide 11: Challenges and Solutions
slide = presentation.slides.add_slide(presentation.slide_layouts[1])
title = slide.shapes.title
content = slide.placeholders[1]
title.text = "Challenges and Solutions"
content.text = (
    "Challenges:\n- Requirement changes.\n- Testing delays.\n\n"
    "Solutions:\n- Agile methods, automated testing."
)

# Slide 12: Conclusion
slide = presentation.slides.add_slide(presentation.slide_layouts[1])
title = slide.shapes.title
content = slide.placeholders[1]
title.text = "Conclusion"
content.text = (
    "Summary:\n- The 'Music Hub' project demonstrates a systematic approach to SDLC.\n"
    "- Ensures quality, efficiency, and user satisfaction.\n\n"
    "Next Steps:\n- Implement further features based on market trends."
)

# Slide 13: References (Optional)
slide = presentation.slides.add_slide(presentation.slide_layouts[1])
title = slide.shapes.title
content = slide.placeholders[1]
title.text = "References"
content.text = "List of tools, frameworks, or methodologies used."

# Save the presentation
file_path = "/mnt/data/Music_Hub_SDLC_Project.pptx"
presentation.save(file_path)

file_path
