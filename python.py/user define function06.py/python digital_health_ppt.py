from pptx import Presentation
from pptx.util import Inches

# Create a presentation
prs = Presentation()

# Define slide titles and contents
slides_content = [
    ("Ethical and Legal Issues in Digital Health Applied to Nursing", "A Comprehensive Overview\nPresented by: [Your Name]\nDate: [Presentation Date]"),
    ("Introduction", "What is Digital Health?\n- Integration of technology into healthcare\n- Includes EHRs, telehealth, AI, and wearable devices\n\nImportance in Nursing\n- Improves efficiency and accessibility\n- Enhances patient safety and communication\n\nObjective\n- Discuss ethical and legal concerns in digital health"),
    ("Overview of Digital Health in Nursing", "- Electronic Health Records (EHRs)\n- Telemedicine & Telehealth\n- Artificial Intelligence in Nursing\n- Wearable Health Technology\n- Mobile Health Apps"),
    ("Ethical Issues in Digital Health", ""),
    ("Patient Confidentiality & Data Privacy", "Importance of HIPAA and GDPR compliance\nRisks of data breaches and cybersecurity threats\nCase studies on data privacy violations"),
    ("Informed Consent in Digital Health", "Challenges in obtaining informed consent\nEnsuring patient autonomy in digital health"),
    ("Digital Divide and Health Equity", "Access disparities in telehealth\nChallenges for elderly and low-income patients"),
    ("AI and Automation in Nursing", "Ethical dilemmas in AI-assisted decision-making\nRisks of bias in AI algorithms"),
    ("Professional Boundaries in Telemedicine", "Challenges in maintaining nurse-patient relationships online\nEthical considerations in virtual consultations"),
    ("Legal Issues in Digital Health", ""),
    ("Compliance with Healthcare Laws", "HIPAA (Health Insurance Portability and Accountability Act)\nGDPR (General Data Protection Regulation)"),
    ("Liability in Telehealth & AI Decisions", "Who is responsible for AI-based nursing recommendations?\nLegal implications of misdiagnosis via telehealth"),
    ("Cybersecurity and Legal Consequences", "Legal actions against data breaches in hospitals\nBest practices for digital security in nursing"),
    ("Documentation and Record-Keeping", "Legal risks of incomplete or incorrect digital records\nStandards for maintaining electronic health records"),
    ("Licensing and Cross-Border Telehealth Challenges", "Legal issues with providing nursing care across state/country lines\nJurisdictional challenges in telemedicine"),
    ("Case Studies & Real-World Examples", "Case study on HIPAA violation and its consequences\nExample of AI misdiagnosis and its legal implications\nSuccess story of ethical telemedicine implementation"),
    ("Best Practices for Nurses in Digital Health", "Adhering to legal and ethical guidelines\nEnsuring data security and patient privacy\nEthical decision-making framework in digital nursing\nStaying updated with digital health regulations"),
    ("Conclusion & Future Trends", "Summary of key takeaways\nFuture trends in digital health and nursing ethics\nCall to action for nurses to uphold ethical and legal standards"),
    ("References & Q&A Slide", "List of sources used\nInvite questions from the audience")
]

# Add slides to presentation
for title, content in slides_content:
    slide_layout = prs.slide_layouts[1]  # Title and content layout
    slide = prs.slides.add_slide(slide_layout)
    title_placeholder = slide.shapes.title
    content_placeholder = slide.placeholders[1]
    title_placeholder.text = title
    content_placeholder.text = content

# Save the presentation
pptx_path = "/mnt/data/Digital_Health_Nursing.pptx"
prs.save(pptx_path)
pptx_path
